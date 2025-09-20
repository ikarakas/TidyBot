from typing import Dict, Any, List, Optional
from pathlib import Path
import logging
from datetime import datetime
import re
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation
import chardet
from collections import Counter
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
import asyncio

logger = logging.getLogger(__name__)

try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')
    nltk.download('stopwords')


class DocumentAnalyzer:
    def __init__(self):
        self.stop_words = set(stopwords.words('english'))
        self.date_patterns = [
            r'\d{4}-\d{2}-\d{2}',
            r'\d{2}/\d{2}/\d{4}',
            r'\d{2}-\d{2}-\d{4}',
            r'\d{1,2}\s+(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4}',
        ]
    
    async def analyze(self, file_path: Path) -> Dict[str, Any]:
        extension = file_path.suffix.lower()
        
        analyzers = {
            '.pdf': self._analyze_pdf,
            '.docx': self._analyze_docx,
            '.doc': self._analyze_docx,
            '.txt': self._analyze_text,
            '.xlsx': self._analyze_excel,
            '.xls': self._analyze_excel,
            '.pptx': self._analyze_powerpoint,
            '.ppt': self._analyze_powerpoint,
        }
        
        analyzer = analyzers.get(extension, self._analyze_text)
        
        try:
            result = await analyzer(file_path)
            result['extension'] = extension
            result['type'] = 'document'
            return result
        except Exception as e:
            logger.error(f"Error analyzing document {file_path}: {e}")
            return {
                'error': str(e),
                'type': 'document',
                'extension': extension
            }
    
    async def _analyze_pdf(self, file_path: Path) -> Dict[str, Any]:
        analysis = {
            'format': 'PDF',
            'pages': 0,
            'text': '',
            'metadata': {},
            'tables_detected': False,
            'images_detected': False
        }
        
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                analysis['pages'] = len(reader.pages)
                
                if reader.metadata:
                    analysis['metadata'] = {
                        'title': reader.metadata.get('/Title', ''),
                        'author': reader.metadata.get('/Author', ''),
                        'subject': reader.metadata.get('/Subject', ''),
                        'creator': reader.metadata.get('/Creator', ''),
                        'creation_date': str(reader.metadata.get('/CreationDate', '')),
                        'modification_date': str(reader.metadata.get('/ModDate', ''))
                    }
                
                text_parts = []
                for i, page in enumerate(reader.pages[:10]):
                    text = page.extract_text()
                    text_parts.append(text)
                    
                    if '/XObject' in page.get('/Resources', {}).get('/XObject', {}):
                        analysis['images_detected'] = True
                
                analysis['text'] = '\n'.join(text_parts)
                
                analysis.update(self._analyze_text_content(analysis['text']))
                
        except Exception as e:
            logger.error(f"PDF analysis error: {e}")
            analysis['error'] = str(e)
        
        return analysis
    
    async def _analyze_docx(self, file_path: Path) -> Dict[str, Any]:
        analysis = {
            'format': 'DOCX',
            'paragraphs': 0,
            'tables': 0,
            'images': 0,
            'text': '',
            'metadata': {}
        }
        
        try:
            doc = Document(file_path)
            
            analysis['paragraphs'] = len(doc.paragraphs)
            analysis['tables'] = len(doc.tables)
            
            if doc.core_properties:
                analysis['metadata'] = {
                    'title': doc.core_properties.title or '',
                    'author': doc.core_properties.author or '',
                    'subject': doc.core_properties.subject or '',
                    'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                    'modified': str(doc.core_properties.modified) if doc.core_properties.modified else '',
                    'revision': doc.core_properties.revision or 0
                }
            
            text_parts = []
            for paragraph in doc.paragraphs[:100]:
                text_parts.append(paragraph.text)
            
            analysis['text'] = '\n'.join(text_parts)
            
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    analysis['images'] += 1
            
            analysis.update(self._analyze_text_content(analysis['text']))
            
        except Exception as e:
            logger.error(f"DOCX analysis error: {e}")
            analysis['error'] = str(e)
        
        return analysis
    
    async def _analyze_excel(self, file_path: Path) -> Dict[str, Any]:
        analysis = {
            'format': 'Excel',
            'sheets': [],
            'total_rows': 0,
            'total_columns': 0,
            'has_formulas': False,
            'has_charts': False,
            'sample_data': []
        }
        
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_info = {
                    'name': sheet_name,
                    'rows': sheet.max_row,
                    'columns': sheet.max_column
                }
                
                analysis['sheets'].append(sheet_info)
                analysis['total_rows'] += sheet.max_row
                analysis['total_columns'] = max(analysis['total_columns'], sheet.max_column)
                
                sample_rows = []
                for row in sheet.iter_rows(max_row=5, values_only=True):
                    sample_rows.append([str(cell)[:50] if cell else '' for cell in row[:5]])
                
                if sample_rows:
                    analysis['sample_data'].append({
                        'sheet': sheet_name,
                        'data': sample_rows
                    })
            
            workbook.close()
            
            analysis['summary'] = self._generate_excel_summary(analysis)
            
        except Exception as e:
            logger.error(f"Excel analysis error: {e}")
            analysis['error'] = str(e)
        
        return analysis
    
    async def _analyze_powerpoint(self, file_path: Path) -> Dict[str, Any]:
        analysis = {
            'format': 'PowerPoint',
            'slides': 0,
            'titles': [],
            'text': '',
            'has_images': False,
            'has_charts': False,
            'metadata': {}
        }
        
        try:
            presentation = Presentation(file_path)
            analysis['slides'] = len(presentation.slides)
            
            if presentation.core_properties:
                analysis['metadata'] = {
                    'title': presentation.core_properties.title or '',
                    'author': presentation.core_properties.author or '',
                    'subject': presentation.core_properties.subject or '',
                    'created': str(presentation.core_properties.created) if presentation.core_properties.created else '',
                    'modified': str(presentation.core_properties.modified) if presentation.core_properties.modified else ''
                }
            
            text_parts = []
            for i, slide in enumerate(presentation.slides[:20]):
                if slide.shapes.title:
                    title = slide.shapes.title.text
                    analysis['titles'].append(title)
                    text_parts.append(title)
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
                    
                    if shape.shape_type == 13:
                        analysis['has_images'] = True
                    elif shape.shape_type == 3:
                        analysis['has_charts'] = True
            
            analysis['text'] = '\n'.join(text_parts)
            analysis.update(self._analyze_text_content(analysis['text']))
            
        except Exception as e:
            logger.error(f"PowerPoint analysis error: {e}")
            analysis['error'] = str(e)
        
        return analysis
    
    async def _analyze_text(self, file_path: Path) -> Dict[str, Any]:
        analysis = {
            'format': 'Text',
            'encoding': 'unknown',
            'lines': 0,
            'words': 0,
            'characters': 0,
            'text': ''
        }
        
        try:
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                detected = chardet.detect(raw_data)
                encoding = detected['encoding'] or 'utf-8'
                analysis['encoding'] = encoding
                
                text = raw_data.decode(encoding, errors='ignore')
                analysis['text'] = text[:10000]
                analysis['lines'] = len(text.splitlines())
                analysis['words'] = len(text.split())
                analysis['characters'] = len(text)
                
                analysis.update(self._analyze_text_content(text))
                
        except Exception as e:
            logger.error(f"Text analysis error: {e}")
            analysis['error'] = str(e)
        
        return analysis
    
    def _analyze_text_content(self, text: str) -> Dict[str, Any]:
        if not text:
            return {}
        
        analysis = {
            'language': self._detect_language(text),
            'keywords': self._extract_keywords(text),
            'dates': self._extract_dates(text),
            'emails': self._extract_emails(text),
            'urls': self._extract_urls(text),
            'summary': self._generate_summary(text)
        }
        
        return analysis
    
    def _detect_language(self, text: str) -> str:
        try:
            from langdetect import detect
            return detect(text[:500])
        except:
            return 'en'
    
    def _extract_keywords(self, text: str, num_keywords: int = 10) -> List[str]:
        try:
            tokens = word_tokenize(text.lower())
            
            words = [token for token in tokens if token.isalnum() and token not in self.stop_words and len(token) > 2]
            
            word_freq = Counter(words)
            
            return [word for word, _ in word_freq.most_common(num_keywords)]
        except:
            return []
    
    def _extract_dates(self, text: str) -> List[str]:
        dates = []
        for pattern in self.date_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            dates.extend(matches)
        return list(set(dates))[:5]
    
    def _extract_emails(self, text: str) -> List[str]:
        email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
        emails = re.findall(email_pattern, text)
        return list(set(emails))[:5]
    
    def _extract_urls(self, text: str) -> List[str]:
        url_pattern = r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
        urls = re.findall(url_pattern, text)
        return list(set(urls))[:5]
    
    def _generate_summary(self, text: str, max_length: int = 200) -> str:
        if len(text) <= max_length:
            return text
        
        sentences = text.split('.')
        summary = []
        current_length = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if sentence and current_length + len(sentence) <= max_length:
                summary.append(sentence)
                current_length += len(sentence)
            elif current_length > 0:
                break
        
        return '. '.join(summary) + '...' if summary else text[:max_length] + '...'
    
    def _generate_excel_summary(self, analysis: Dict[str, Any]) -> str:
        parts = []
        
        if analysis['sheets']:
            parts.append(f"{len(analysis['sheets'])} sheets")
        
        if analysis['total_rows']:
            parts.append(f"{analysis['total_rows']} total rows")
        
        if analysis['total_columns']:
            parts.append(f"up to {analysis['total_columns']} columns")
        
        return ', '.join(parts) if parts else 'Empty spreadsheet'